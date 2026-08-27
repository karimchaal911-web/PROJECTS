"""Build the PDF and PPTX fallbacks from the captured presentation frames.

The fallbacks are not a separate deck. They are the same experience, frame for
frame, so a projector failure costs the presenter production value but never
content: same scenes, same order, same numbers, same speaker notes.

  python final_presentation_claude/data_build/build_fallbacks.py
"""

from __future__ import annotations

import json
import re
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "exports" / "screenshots"
NOTES = ROOT / "speaker_notes" / "quick_cues.md"
PDF_OUT = ROOT / "FINAL_MAP_Soluble_Digitalization_Soutenance_Claude.pdf"
PPTX_OUT = ROOT / "FINAL_MAP_Soluble_Digitalization_Soutenance_Claude.pptx"

W_IN, H_IN = 13.333, 7.5  # 16:9


def load_frames():
    files = sorted(SHOTS.glob("step-*.png"))
    if not files:
        raise SystemExit("No captures in exports/screenshots — run npm run capture first.")
    manifest = {}
    mf = SHOTS / "manifest.json"
    if mf.exists():
        manifest = {s["file"]: s for s in json.loads(mf.read_text()).get("steps", [])}
    return files, manifest


def parse_cues() -> dict:
    """Pull the per-step bullet lists out of quick_cues.md by its ### headings."""
    if not NOTES.exists():
        return {}
    out, key, buf = {}, None, []
    for line in NOTES.read_text(encoding="utf-8").splitlines():
        m = re.match(r"^### (\S+) · (.+)$", line.strip())
        if m:
            if key:
                out[key] = (title, [b for b in buf if b])
            key, title, buf = m.group(1), m.group(2), []
            continue
        if key and line.strip().startswith("- "):
            buf.append(line.strip()[2:])
        elif key and line.startswith("#"):
            out[key] = (title, [b for b in buf if b])
            key, buf = None, []
    if key:
        out[key] = (title, [b for b in buf if b])
    return out


def cue_key(index: int, scene: str, beat: int, beats_in_scene: int) -> str:
    """quick_cues.md keys are '05a', '05b', … for multi-beat scenes, '12' otherwise."""
    if beats_in_scene > 1:
        return f"{scene}{chr(ord('a') + beat - 1)}"
    return scene


def strip_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


def build_pdf(files):
    doc = fitz.open()
    for f in files:
        with Image.open(f) as im:
            w, h = im.size
        # 16:9 at 96 dpi-equivalent points
        page = doc.new_page(width=W_IN * 72, height=H_IN * 72)
        page.insert_image(fitz.Rect(0, 0, W_IN * 72, H_IN * 72), filename=str(f))
    doc.set_metadata({
        "title": "Intelligent Digitalization of Soluble MAP Production",
        "author": "OCP · ENSAM — soutenance",
        "subject": "Static fallback for the interactive 3D soutenance",
        "keywords": "soluble MAP, rotary dryer, soft sensor, One-Class SVM, Power BI",
    })
    doc.save(PDF_OUT, deflate=True, garbage=4)
    doc.close()
    print("wrote %s  (%d pages, %.1f MB)"
          % (PDF_OUT.name, len(files), PDF_OUT.stat().st_size / 1e6))


def build_pptx(files, manifest, cues):
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]

    # How many beats each scene has, so cue keys can be derived.
    scene_beats = {}
    for f in files:
        m = re.search(r"scene-(\d+)_beat-(\d+)", f.name)
        if m:
            scene_beats[m.group(1)] = max(scene_beats.get(m.group(1), 0), int(m.group(2)))

    for i, f in enumerate(files):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(f), 0, 0, width=prs.slide_width, height=prs.slide_height)

        m = re.search(r"scene-(\d+)_beat-(\d+)", f.name)
        scene, beat = (m.group(1), int(m.group(2))) if m else ("00", 1)
        key = cue_key(i, scene, beat, scene_beats.get(scene, 1))
        title, bullets = cues.get(key, (manifest.get(f.name, {}).get("hero", ""), []))

        notes = slide.notes_slide.notes_text_frame
        notes.text = f"SCENE {scene} · {title}"
        for b in bullets:
            p = notes.add_paragraph()
            p.text = strip_md(b)
            p.font.size = Pt(12)

    prs.core_properties.title = "Intelligent Digitalization of Soluble MAP Production"
    prs.core_properties.author = "OCP · ENSAM — soutenance"
    prs.core_properties.comments = (
        "Fallback deck. The primary deliverable is the interactive React + Three.js "
        "presentation in final_presentation_claude/web. Speaker notes on each slide "
        "are the same cues as speaker_notes/quick_cues.md."
    )
    prs.save(PPTX_OUT)
    print("wrote %s  (%d slides, %.1f MB)"
          % (PPTX_OUT.name, len(files), PPTX_OUT.stat().st_size / 1e6))


def main():
    files, manifest = load_frames()
    cues = parse_cues()
    print("frames: %d · cue cards parsed: %d" % (len(files), len(cues)))
    build_pdf(files)
    build_pptx(files, manifest, cues)


if __name__ == "__main__":
    main()
